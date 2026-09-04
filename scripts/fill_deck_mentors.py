#!/usr/bin/env python3
"""Fill the mandatory mentor rows on the team slide, and fix a registration number.

Both mentor rows on slide 7 are marked *(Mandatory)* and both are empty, which
is a hard submission blocker independent of any code. This edits them in place
rather than by hand so the change is repeatable and so nothing else on the
slide moves.

    ~/carla-venv/bin/python3 scripts/fill_deck_mentors.py \\
        --faculty "DR A B SHARMA|CINTEL|-|-|-|ab.sharma@srmist.edu.in|9876543210" \\
        --industry "MS C D RAO|Acme Robotics|-|-|-|cd.rao@acme.com|9876543211" \\
        --fix-reg "SWETANK KUMAR=RA2511026010134"

Fields are pipe-separated in the table's own column order:
``Name|Dept|Year|Semester|Gender|Email|Mobile``. Anything you leave empty stays
empty; a mentor has no year of study, so ``-`` is the honest filler.

**The deck is deliberately gitignored** -- it carries six people's names,
registration numbers, emails and mobile numbers and the repository is public.
Do not commit it, and make the repo private before you ever do.
"""

from __future__ import annotations

import argparse
import shutil
import sys
from datetime import datetime
from pathlib import Path

DECK = Path("deck/SRMIST SIH2026-TEAM DIVAS.pptx")
#: Table column order on slide 7, after the leading role column.
COLUMNS = ("Name", "Reg. No.", "Dept.", "Year of Study", "Semester", "Gender",
           "Email id", "Mobile no.")


def set_cell(cell, text: str) -> None:
    """Write text while keeping the cell's existing formatting.

    Assigning to ``cell.text`` throws away the run properties and the row
    then renders in a different font from the six rows above it, which on a
    submitted slide looks like a late patch -- because it is one. Writing into
    the first run keeps the font, size and colour the template set.
    """
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra.text = ""
    else:
        para.add_run().text = text


def find_team_table(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.table.rows[0].cells[1].text.strip() == "Name":
                return shape.table
    raise SystemExit("no team table found -- has the deck been restructured?")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--deck", default=str(DECK))
    ap.add_argument("--faculty", default=None,
                    help="Name|Dept|Year|Sem|Gender|Email|Mobile")
    ap.add_argument("--industry", default=None, help="same field order")
    ap.add_argument("--fix-reg", action="append", default=[],
                    metavar="NAME=REGNO",
                    help="correct one member's registration number; repeatable")
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    try:
        from pptx import Presentation
    except ImportError:
        print("needs python-pptx:  ~/carla-venv/bin/pip install python-pptx",
              file=sys.stderr)
        return 2

    deck = Path(args.deck)
    if not deck.exists():
        print(f"no deck at {deck}", file=sys.stderr)
        return 2

    prs = Presentation(str(deck))
    table = find_team_table(prs)

    def fill(role_fragment: str, spec: str) -> None:
        fields = [f.strip() for f in spec.split("|")]
        if len(fields) != 7:
            raise SystemExit(
                f"--{role_fragment} needs 7 pipe-separated fields "
                f"(Name|Dept|Year|Sem|Gender|Email|Mobile), got {len(fields)}"
            )
        # Reg. No. is column 2 and a mentor has none, so the fields skip it.
        values = [fields[0], ""] + fields[1:]
        for row in table.rows:
            if role_fragment in row.cells[0].text.lower():
                for col, value in enumerate(values, start=1):
                    if value:
                        set_cell(row.cells[col], value)
                print(f"  {row.cells[0].text.strip()}: {fields[0]}")
                return
        raise SystemExit(f"no row matching {role_fragment!r}")

    if args.faculty:
        fill("faculty", args.faculty)
    if args.industry:
        fill("industry", args.industry)

    for item in args.fix_reg:
        name, _, regno = item.partition("=")
        for row in table.rows:
            if row.cells[1].text.strip().lower() == name.strip().lower():
                old = row.cells[2].text.strip()
                set_cell(row.cells[2], regno.strip())
                print(f"  {name.strip()}: reg. no. {old} -> {regno.strip()} "
                      f"({len(regno.strip())} chars)")
                break
        else:
            raise SystemExit(f"no member named {name!r} on the slide")

    # Report anything still blocking submission, because a silent success here
    # is worse than useless -- the whole point of the script is the blocker.
    outstanding = [
        row.cells[0].text.strip()
        for row in table.rows
        if "mandatory" in row.cells[0].text.lower() and not row.cells[1].text.strip()
    ]
    # Skip the header row, whose "Reg. No." caption is not a registration
    # number and would otherwise be reported as a malformed one forever.
    lengths = {row.cells[1].text.strip(): row.cells[2].text.strip()
               for row in list(table.rows)[1:] if row.cells[2].text.strip()}
    odd = {n: r for n, r in lengths.items() if len(r) != 15}

    if args.dry_run:
        print("\n--dry-run: nothing written")
    else:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        backup = deck.with_name(f"{deck.stem}.backup-{stamp}{deck.suffix}")
        shutil.copy(deck, backup)
        prs.save(str(deck))
        print(f"\nsaved {deck}\nbackup {backup}")

    if outstanding:
        print(f"\nSTILL BLOCKING SUBMISSION: {outstanding} empty and mandatory")
    if odd:
        print(f"malformed registration numbers (expected 15 chars): {odd}")
    if not outstanding and not odd:
        print("\nno submission blockers left on this slide")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
